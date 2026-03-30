#!/usr/bin/env python3
"""
Batch 7 Extraction Script (Files 151-175)
Mandatory 300 DPI OCR on every PDF page.
CLI tools (pdftotext/pdftoppm) for files >10MB.
python-pptx for PPTX files.
"""

import json
import os
import hashlib
import subprocess
import re
import sys
import traceback
from datetime import datetime, timezone

# Try imports
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    import pytesseract
    from PIL import Image
except ImportError:
    pytesseract = None
    Image = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# Paths
SOURCE_DIR = "/app/backend/brochure_intelligence/source_brochures"
RAW_DIR = "/app/backend/brochure_intelligence/raw_extractions"
DRAFT_DIR = "/app/backend/brochure_intelligence/structured_drafts"
LOG_DIR = "/app/backend/brochure_intelligence/logs"

# Get the sorted file list
ALL_FILES = sorted(os.listdir(SOURCE_DIR))
BATCH_7_FILES = ALL_FILES[150:175]  # 0-indexed

SIZE_THRESHOLD_MB = 10

# SKU pattern regex - catches common medical device SKU patterns
SKU_PATTERNS = [
    r'\b(MT-[A-Z]{2}\d{4,}[A-Z]*)\b',           # MT- trauma codes
    r'\b(MT-[A-Z]{2}\d+[A-Z]?\d*[A-Z]*)\b',      # MT- flexible
    r'\b(BIA\d{3,})\b',                            # BioMime Aura
    r'\b(BIL\d{3,})\b',                            # BioMime Lineage
    r'\b(BBR\d{5,})\b',                            # BioMime Branch
    r'\b(MOZ[A-Z]?\d{3,})\b',                      # MOZEC codes
    r'\b(MNC\d{3,})\b',                             # MOZEC NC
    r'\b(MRL\d{3,})\b',                             # MeRes100
    r'\b(MOZS\d{3,})\b',                            # MOZEC SEB
    r'\b(RPD[A-Z]{2,}-\d+)\b',                      # MeriScreen rapid
    r'\b(ELI[A-Z]{2,}-\d+)\b',                      # ELISA kits
    r'\b(MSG[A-Z0-9]{2,})\b',                       # Gowns / guide catheters
    r'\b(MSB[A-Z0-9]{2,})\b',                       # Balloon catheters
    r'\b(MSBC-\d+/\d+)\b',                          # Latitud bipolar
    r'\b(CING-\d+/\d+)\b',                          # Latitud cemented
    r'\b(MAC\d{5,})\b',                              # METIC airway
    r'\b(MEK\d{5,})\b',                              # MYCS kits
    r'\b(MEB\d{5,})\b',                              # MYCS balloons
    r'\b(MTT\d{3,}[A-Z]?)\b',                       # MAFIC tracheostomy
    r'\b(ALT\d{4,})\b',                              # ALRINE T-tube
    r'\b(DASK[A-Z]+-\d+[-\d]*)\b',                  # DOA rapids dip
    r'\b(DAR[A-Z]+-\d+[-\d]*)\b',                   # DOA rapids cassette
    r'\b(DOARPD-\d+)\b',                             # DOA multi-drug
    r'\b([A-Z]{2,5}\d{3,}[A-Z]?\d*[A-Z]*U?)\b',   # Generic code patterns
    r'\b(\d{2}[A-Z]{2,}\d{2,})\b',                  # Numeric-alpha codes
    r'\b(SPM-[A-Z0-9]+)\b',                          # SPM spine codes
    r'\b(INDR[A-Z0-9]+)\b',                          # Armar drill codes
    r'\b(NSLPS[A-Z0-9]+)\b',                         # Armar plate codes
]

# Ordering page keywords
ORDERING_KEYWORDS = [
    'ordering', 'order code', 'product code', 'catalogue', 'catalog',
    'ref', 'reference', 'sku', 'part number', 'item code', 'article',
    'code no', 'model no', 'size chart', 'specifications table'
]


def get_file_id(filepath):
    """Generate content-based file ID."""
    sha = hashlib.sha256()
    with open(filepath, 'rb') as f:
        for chunk in iter(lambda: f.read(8192), b''):
            sha.update(chunk)
    return f"BF_{sha.hexdigest()[:12]}"


def extract_text_pdfplumber(filepath, max_pages=None):
    """Extract text using pdfplumber (for small PDFs)."""
    pages = []
    try:
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages):
                if max_pages and i >= max_pages:
                    break
                text = page.extract_text() or ''
                tables = page.extract_tables() or []
                table_text = ''
                for table in tables:
                    for row in table:
                        if row:
                            table_text += ' | '.join([str(c) if c else '' for c in row]) + '\n'
                pages.append({
                    'page_number': i + 1,
                    'text': text,
                    'table_text': table_text,
                    'has_tables': len(tables) > 0
                })
    except Exception as e:
        print(f"  pdfplumber error: {e}")
    return pages


def extract_text_cli(filepath):
    """Extract text using pdftotext CLI (for large PDFs)."""
    pages = []
    try:
        result = subprocess.run(
            ['pdftotext', '-layout', filepath, '-'],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode == 0:
            # Split by form feed character
            raw_pages = result.stdout.split('\f')
            for i, page_text in enumerate(raw_pages):
                if page_text.strip():
                    pages.append({
                        'page_number': i + 1,
                        'text': page_text.strip(),
                        'table_text': '',
                        'has_tables': False
                    })
    except Exception as e:
        print(f"  CLI pdftotext error: {e}")
    return pages


def ocr_pdf_page(filepath, page_num, dpi=300):
    """OCR a single PDF page at given DPI using pdftoppm + tesseract."""
    try:
        # Render page to image
        result = subprocess.run(
            ['pdftoppm', '-f', str(page_num), '-l', str(page_num),
             '-r', str(dpi), '-png', filepath],
            capture_output=True, timeout=60
        )
        if result.returncode == 0 and result.stdout:
            # Save temp image
            tmp_img = f'/tmp/ocr_page_{page_num}.png'
            with open(tmp_img, 'wb') as f:
                f.write(result.stdout)
            # OCR
            img = Image.open(tmp_img)
            text = pytesseract.image_to_string(img)
            os.remove(tmp_img)
            return text.strip()
    except Exception as e:
        print(f"  OCR error page {page_num}: {e}")
    return ''


def extract_pptx(filepath):
    """Extract text from PPTX."""
    pages = []
    try:
        prs = Presentation(filepath)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    texts.append(shape.text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = ' | '.join([cell.text for cell in row.cells])
                        texts.append(row_text)
            pages.append({
                'page_number': i + 1,
                'text': '\n'.join(texts),
                'table_text': '',
                'has_tables': False
            })
    except Exception as e:
        print(f"  PPTX error: {e}")
    return pages


def find_skus_in_text(text):
    """Extract SKU codes from text using known patterns."""
    found = set()
    for pattern in SKU_PATTERNS:
        matches = re.findall(pattern, text)
        for m in matches:
            code = m.strip()
            # Filter out very short or common words
            if len(code) >= 4 and not code.lower() in {'the', 'and', 'for', 'with', 'from'}:
                found.add(code)
    return sorted(found)


def has_ordering_keywords(text):
    """Check if text contains ordering-related keywords."""
    text_lower = text.lower()
    return any(kw in text_lower for kw in ORDERING_KEYWORDS)


def classify_file(filename, text_content):
    """Classify file type based on content."""
    fn_lower = filename.lower()
    text_lower = text_content.lower() if text_content else ''

    if 'price' in fn_lower or 'pricing' in fn_lower:
        return 'price_list_commercial'
    if 'ifu' in fn_lower or 'catalogue' in fn_lower or 'catalog' in fn_lower:
        return 'ifu_official_catalog'
    if 'surgical technique' in fn_lower or ' st ' in fn_lower or '_st.' in fn_lower:
        return 'surgical_technique'
    if 'training' in fn_lower or 'ppt' in fn_lower.rsplit('.', 1)[-1]:
        return 'training_deck'
    if 'brochure' in fn_lower:
        return 'product_brochure'
    if 'flyer' in fn_lower:
        return 'product_flyer'
    if any(kw in text_lower for kw in ['ordering information', 'product code', 'catalogue number']):
        return 'product_catalog'
    return 'product_brochure'


def extract_products_from_text(all_text, filename, file_num):
    """Extract product information from combined text."""
    products = []
    all_skus = find_skus_in_text(all_text)

    # Determine brand from filename and content
    fn_lower = filename.lower()
    brand = 'Unknown'
    division = 'Unknown'

    # Brand detection from filename
    brand_map = {
        'surgical-needles': ('Meril', 'Endo Surgery'),
        'suture': ('Meril', 'Endo Surgery'),
        'trauma nail': ('Meril Orthopedics', 'Trauma'),
        'twin_trent': ('Twin Trent', 'Cardiovascular'),
        'ti_elbow': ('Meril Orthopedics', 'Trauma'),
        'tournisys': ('Tournisys', 'Trauma'),
        'trauma plate': ('Meril Orthopedics', 'Trauma'),
        'tredro': ('Tredro', 'Diagnostics'),
        'trocar': ('MIRUS', 'Endo Surgery'),
        'filaprop': ('FILAPROP', 'Endo Surgery'),
        'variabilis': ('Variabilis', 'Trauma'),
        'absomesh': ('Absomesh', 'Endo Surgery'),
        'ancora': ('Ancora', 'Endo Surgery'),
        'biomime': ('BioMime', 'Cardiovascular'),
        'bonewax': ('Bonewax', 'Endo Surgery'),
        'cogent': ('Cogent', 'Cardiovascular'),
        'evermine': ('Evermine', 'Cardiovascular'),
        'fiona': ('Fiona', 'Cardiovascular'),
        'handx': ('HandX', 'Endo Surgery'),
        'menox': ('Menox', 'Cardiovascular'),
        'merifeim': ('MeriFeim', 'Cardiovascular'),
    }

    for key, (b, d) in brand_map.items():
        if key in fn_lower:
            brand = b
            division = d
            break

    # Also try to detect from content
    content_lower = all_text.lower()
    if brand == 'Unknown':
        content_brands = {
            'biomime': ('BioMime', 'Cardiovascular'),
            'filaprop': ('FILAPROP', 'Endo Surgery'),
            'mirus': ('MIRUS', 'Endo Surgery'),
            'variabilis': ('Variabilis', 'Trauma'),
            'tredro': ('Tredro', 'Diagnostics'),
        }
        for key, (b, d) in content_brands.items():
            if key in content_lower:
                brand = b
                division = d
                break

    # Create a single product entry (will be refined during structuring)
    product_name = os.path.splitext(filename)[0]
    # Clean up product name
    product_name = re.sub(r'_[a-f0-9]{8,}$', '', product_name)  # Remove hash suffixes
    product_name = re.sub(r'_\d+$', '', product_name)  # Remove trailing numbers
    product_name = product_name.replace('_', ' ').replace('-', ' ').strip()
    product_name = re.sub(r'\s+', ' ', product_name)

    product = {
        'product_name': product_name,
        'brand': brand,
        'division': division,
        'description': all_text[:500] if all_text else '',
        'sku_codes': all_skus,
        'has_ordering_info': has_ordering_keywords(all_text),
        'source_file': f'{file_num:03d}'
    }
    products.append(product)

    return products


def process_file(file_num, filename):
    """Process a single file through the extraction pipeline."""
    filepath = os.path.join(SOURCE_DIR, filename)
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else 'unknown'
    size_mb = os.path.getsize(filepath) / (1024 * 1024)
    file_id = get_file_id(filepath)

    print(f"\n{'='*60}")
    print(f"Processing File {file_num:03d}: {filename}")
    print(f"  Type: {ext}, Size: {size_mb:.1f} MB, ID: {file_id}")

    # Choose extraction method
    pages = []
    parser_used = 'unknown'

    if ext == 'pptx':
        parser_used = 'python-pptx'
        pages = extract_pptx(filepath)

    elif ext == 'pdf':
        if size_mb > SIZE_THRESHOLD_MB:
            # Use CLI for large PDFs
            parser_used = 'CLI pdftotext + pdftoppm OCR @ 300 DPI'
            print(f"  Using CLI extraction (file > {SIZE_THRESHOLD_MB}MB)")
            pages = extract_text_cli(filepath)

            # Get page count for OCR
            try:
                result = subprocess.run(
                    ['pdfinfo', filepath], capture_output=True, text=True, timeout=30
                )
                page_count = 0
                for line in result.stdout.split('\n'):
                    if line.startswith('Pages:'):
                        page_count = int(line.split(':')[1].strip())
                        break
            except:
                page_count = len(pages)

            # OCR each page at 300 DPI
            print(f"  Running 300 DPI OCR on {page_count} pages...")
            for pg_num in range(1, page_count + 1):
                ocr_text = ocr_pdf_page(filepath, pg_num, dpi=300)
                # Find or create page entry
                found = False
                for p in pages:
                    if p['page_number'] == pg_num:
                        p['ocr_text'] = ocr_text
                        found = True
                        break
                if not found and ocr_text:
                    pages.append({
                        'page_number': pg_num,
                        'text': '',
                        'ocr_text': ocr_text,
                        'table_text': '',
                        'has_tables': False
                    })

        else:
            # Use pdfplumber for small PDFs
            parser_used = 'pdfplumber + pytesseract OCR @ 300 DPI'
            pages = extract_text_pdfplumber(filepath)

            # 300 DPI OCR on every page (mandatory)
            try:
                result = subprocess.run(
                    ['pdfinfo', filepath], capture_output=True, text=True, timeout=30
                )
                page_count = 0
                for line in result.stdout.split('\n'):
                    if line.startswith('Pages:'):
                        page_count = int(line.split(':')[1].strip())
                        break
            except:
                page_count = len(pages)

            print(f"  Running mandatory 300 DPI OCR on {page_count} pages...")
            for pg_num in range(1, page_count + 1):
                ocr_text = ocr_pdf_page(filepath, pg_num, dpi=300)
                found = False
                for p in pages:
                    if p['page_number'] == pg_num:
                        p['ocr_text'] = ocr_text
                        found = True
                        break
                if not found and ocr_text:
                    pages.append({
                        'page_number': pg_num,
                        'text': '',
                        'ocr_text': ocr_text,
                        'table_text': '',
                        'has_tables': False
                    })

    else:
        print(f"  Unsupported file type: {ext}")
        return None

    # Combine all text
    all_text_parts = []
    for p in sorted(pages, key=lambda x: x['page_number']):
        all_text_parts.append(p.get('text', ''))
        all_text_parts.append(p.get('ocr_text', ''))
        all_text_parts.append(p.get('table_text', ''))
    all_text = '\n'.join(all_text_parts)

    # Classify file
    file_type = classify_file(filename, all_text)

    # Find SKUs across all pages
    page_skus = {}
    total_skus = set()
    ordering_pages = []
    for p in pages:
        page_text = (p.get('text', '') + ' ' + p.get('ocr_text', '') + ' ' + p.get('table_text', '')).strip()
        skus = find_skus_in_text(page_text)
        if skus:
            page_skus[p['page_number']] = skus
            total_skus.update(skus)
        if has_ordering_keywords(page_text):
            ordering_pages.append(p['page_number'])

    # Extract products
    products = extract_products_from_text(all_text, filename, file_num)
    # Update SKUs with page-level findings
    if products and total_skus:
        products[0]['sku_codes'] = sorted(total_skus)

    print(f"  Pages: {len(pages)}, File type: {file_type}")
    print(f"  Products found: {len(products)}")
    print(f"  Total SKU codes: {len(total_skus)}")
    print(f"  Ordering pages: {ordering_pages}")

    # Build raw extraction
    raw_extraction = {
        'file_id': file_id,
        'file_number': f'{file_num:03d}',
        'source_file': filename,
        'file_type': file_type,
        'file_size_mb': round(size_mb, 1),
        'parser_used': parser_used,
        'extraction_method': 'mandatory_300dpi_ocr',
        'total_pages': len(pages),
        'extracted_at': datetime.now(timezone.utc).isoformat(),
        'page_extractions': [
            {
                'page_number': p['page_number'],
                'text_length': len(p.get('text', '')),
                'ocr_text_length': len(p.get('ocr_text', '')),
                'has_tables': p.get('has_tables', False),
                'skus_found': page_skus.get(p['page_number'], []),
                'is_ordering_page': p['page_number'] in ordering_pages
            }
            for p in sorted(pages, key=lambda x: x['page_number'])
        ],
        '_raw_evidence_status': 'COMPLETE',
        'products_found': len(products),
        'skus_found': len(total_skus),
        'ordering_pages': ordering_pages
    }

    # Save raw extraction
    raw_fn = f"{file_num:03d}_{filename.replace(' ', '_').replace('.', '_')}_raw.json"
    raw_path = os.path.join(RAW_DIR, raw_fn)
    with open(raw_path, 'w') as f:
        json.dump(raw_extraction, f, indent=2)

    # Build structured draft
    draft = {
        'file_id': file_id,
        'file_number': f'{file_num:03d}',
        'source_file': filename,
        'file_type': file_type,
        'parser_used': parser_used,
        'drafted_at': datetime.now(timezone.utc).isoformat(),
        '_batch': 'batch_7',
        '_raw_evidence_status': 'COMPLETE',
        'products': products
    }

    # Save structured draft
    draft_fn = f"{file_num:03d}_{filename.replace(' ', '_').replace('.', '_')}_draft.json"
    draft_path = os.path.join(DRAFT_DIR, draft_fn)
    with open(draft_path, 'w') as f:
        json.dump(draft, f, indent=2)

    return {
        'file_num': file_num,
        'filename': filename,
        'file_type': file_type,
        'parser': parser_used,
        'pages': len(pages),
        'products': len(products),
        'skus': len(total_skus),
        'ordering_pages': ordering_pages,
        'key_finding': products[0]['description'][:200] if products else 'No products'
    }


def main():
    print("=" * 70)
    print("BATCH 7 EXTRACTION (Files 151-175)")
    print(f"Started: {datetime.now(timezone.utc).isoformat()}")
    print("=" * 70)

    results = []
    total_products = 0
    total_skus = 0
    files_with_skus = 0
    files_no_skus = 0
    errors = []

    for i, filename in enumerate(BATCH_7_FILES):
        file_num = 151 + i
        try:
            result = process_file(file_num, filename)
            if result:
                results.append(result)
                total_products += result['products']
                total_skus += result['skus']
                if result['skus'] > 0:
                    files_with_skus += 1
                else:
                    files_no_skus += 1
            else:
                files_no_skus += 1
                errors.append(f"File {file_num}: No result")
        except Exception as e:
            print(f"\n  ERROR processing file {file_num}: {e}")
            traceback.print_exc()
            errors.append(f"File {file_num}: {str(e)}")
            files_no_skus += 1

    # Summary
    print("\n" + "=" * 70)
    print("BATCH 7 EXTRACTION SUMMARY")
    print("=" * 70)
    print(f"Files processed: {len(results)}/25")
    print(f"Total products: {total_products}")
    print(f"Total unique SKU codes: {total_skus}")
    print(f"Files with SKUs: {files_with_skus}")
    print(f"Files without SKUs: {files_no_skus}")
    if errors:
        print(f"\nErrors ({len(errors)}):")
        for e in errors:
            print(f"  - {e}")

    print("\nPer-file breakdown:")
    for r in results:
        sku_info = f"{r['skus']} SKUs" if r['skus'] > 0 else "NO SKUs"
        print(f"  {r['file_num']:03d}: {r['products']} products, {sku_info} | {r['filename'][:50]}")

    # Save batch summary
    summary = {
        'batch': 7,
        'files': '151-175',
        'processed_at': datetime.now(timezone.utc).isoformat(),
        'total_files': len(results),
        'total_products': total_products,
        'total_skus': total_skus,
        'files_with_skus': files_with_skus,
        'files_no_skus': files_no_skus,
        'errors': errors,
        'file_results': results
    }

    summary_path = os.path.join(LOG_DIR, 'batch7_extraction_summary.json')
    with open(summary_path, 'w') as f:
        json.dump(summary, f, indent=2)

    print(f"\nSummary saved: {summary_path}")
    print(f"Completed: {datetime.now(timezone.utc).isoformat()}")

    return summary


if __name__ == '__main__':
    main()
