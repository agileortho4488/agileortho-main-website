"""
Batch 3 (Files 51-75) — Batch 1-Quality Extraction
Rules:
1. Raw evidence first (every page)
2. Ordering/code page identification
3. Full SKU extraction for files with ordering signals
4. Explicit no-SKU verification
5. Structured drafts only after evidence exists
"""
import json
import re
import os
import sys
import traceback
from pathlib import Path
from datetime import datetime, timezone

import pdfplumber
import pytesseract
from PIL import Image

BASE = Path("/app/backend/brochure_intelligence")
RAW_DIR = BASE / "raw_extractions"
DRAFT_DIR = BASE / "structured_drafts"
LOG_DIR = BASE / "logs"
BROCHURES = Path("/tmp/zoho_brochures")

DPI = 300
ORDERING_KEYWORDS = [
    "ordering information", "order information", "product code",
    "catalogue", "catalog", "cat. no", "cat no", "ref.", "sr. no",
    "mat code", "material code", "product description"
]

# Read file list
with open("/tmp/brochure_list.txt") as f:
    ALL_FILES = [line.strip() for line in f.readlines() if line.strip()]


def get_file_for_eid(eid_num):
    """Get brochure filename for extraction ID (1-indexed into sorted list)."""
    idx = eid_num - 1
    if idx < 0 or idx >= len(ALL_FILES):
        return None
    return ALL_FILES[idx]


def extract_raw_evidence(eid_num):
    """Phase 1: Extract raw text from EVERY page. No interpretation."""
    filename = get_file_for_eid(eid_num)
    if not filename:
        return None

    eid = f"{eid_num:03d}"
    filepath = BROCHURES / filename
    file_type = filepath.suffix.lower()

    print(f"\n{'='*100}")
    print(f"EXTRACTING RAW EVIDENCE: {eid} — {filename}")
    print(f"{'='*100}")

    raw = {
        "extraction_id": eid,
        "file_id": f"BF_{eid}_{re.sub(r'[^a-z0-9]+', '_', filename.lower())[:40]}",
        "source_file": filename,
        "file_type": file_type,
        "total_pages": 0,
        "parser_used": "",
        "_raw_text_by_page": {},
        "_tables_by_page": {},
        "_raw_evidence_status": "MISSING",
        "page_extractions": [],
        "extracted_at": datetime.now(timezone.utc).isoformat()
    }

    if file_type == ".pdf":
        raw = _extract_pdf(filepath, raw)
    elif file_type in (".pptx",):
        raw = _extract_pptx(filepath, raw)
    elif file_type in (".xlsx", ".xls"):
        raw = _extract_excel(filepath, raw)
    elif file_type in (".docx", ".doc"):
        raw = _extract_docx(filepath, raw)
    else:
        raw["_raw_evidence_status"] = "UNREADABLE"
        raw["unreadable_log"] = {"reason": f"Unsupported file type: {file_type}"}

    # Save raw extraction
    slug = re.sub(r'[^a-z0-9]+', '_', filename.lower())[:40]
    out_path = RAW_DIR / f"{eid}_{slug}.json"
    with open(out_path, "w") as f:
        json.dump(raw, f, indent=2)
    print(f"  Saved: {out_path.name}")

    return raw


def _extract_pdf(filepath, raw):
    """Extract PDF using pdfplumber + 300 DPI OCR fallback."""
    try:
        with pdfplumber.open(str(filepath)) as pdf:
            raw["total_pages"] = len(pdf.pages)
            raw["parser_used"] = "pdfplumber+pytesseract_300dpi"

            for pg_idx, page in enumerate(pdf.pages):
                pg_num = pg_idx + 1

                # Direct text extraction
                direct_text = page.extract_text() or ""

                # Table extraction
                tables = page.extract_tables() or []

                # OCR at 300 DPI if direct text is thin
                ocr_text = ""
                extraction_method = "pdfplumber_direct"
                if len(direct_text) < 50:
                    try:
                        img = page.to_image(resolution=DPI)
                        ocr_text = pytesseract.image_to_string(img.original)
                        extraction_method = "pytesseract_OCR_300dpi"
                    except Exception as e:
                        ocr_text = ""
                        extraction_method = f"pdfplumber_direct(ocr_failed:{str(e)[:50]})"

                # Use the richer text
                best_text = ocr_text if len(ocr_text) > len(direct_text) else direct_text

                raw["_raw_text_by_page"][str(pg_num)] = best_text

                if tables:
                    raw["_tables_by_page"][str(pg_num)] = [
                        {"table_index": i, "rows": t} for i, t in enumerate(tables)
                    ]

                raw["page_extractions"].append({
                    "page_number": pg_num,
                    "source_file": filepath.name,
                    "extraction_method": extraction_method,
                    "direct_text_chars": len(direct_text),
                    "ocr_text_chars": len(ocr_text),
                    "tables_found": len(tables),
                    "best_text_chars": len(best_text)
                })

                print(f"  Page {pg_num:>2}/{raw['total_pages']}: {len(best_text):>5} chars ({extraction_method}){' +' + str(len(tables)) + ' tables' if tables else ''}")

            raw["_raw_evidence_status"] = "PRESENT"

    except Exception as e:
        raw["_raw_evidence_status"] = "UNREADABLE"
        raw["unreadable_log"] = {"reason": str(e), "traceback": traceback.format_exc()[:500]}
        print(f"  ERROR: {e}")

    return raw


def _extract_pptx(filepath, raw):
    """Extract PowerPoint text slide by slide."""
    try:
        from pptx import Presentation
        prs = Presentation(str(filepath))
        raw["total_pages"] = len(prs.slides)
        raw["parser_used"] = "python-pptx"

        for slide_idx, slide in enumerate(prs.slides):
            pg_num = slide_idx + 1
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    if table_data:
                        if str(pg_num) not in raw["_tables_by_page"]:
                            raw["_tables_by_page"][str(pg_num)] = []
                        raw["_tables_by_page"][str(pg_num)].append({
                            "table_index": len(raw["_tables_by_page"].get(str(pg_num), [])),
                            "rows": table_data
                        })

            page_text = "\n".join(texts)
            raw["_raw_text_by_page"][str(pg_num)] = page_text

            raw["page_extractions"].append({
                "page_number": pg_num,
                "source_file": filepath.name,
                "extraction_method": "python-pptx",
                "best_text_chars": len(page_text)
            })

            print(f"  Slide {pg_num:>2}/{raw['total_pages']}: {len(page_text):>5} chars")

        raw["_raw_evidence_status"] = "PRESENT"

    except Exception as e:
        raw["_raw_evidence_status"] = "UNREADABLE"
        raw["unreadable_log"] = {"reason": str(e)}
        print(f"  ERROR: {e}")

    return raw


def _extract_excel(filepath, raw):
    """Extract Excel as text."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(filepath), data_only=True)
        raw["total_pages"] = len(wb.sheetnames)
        raw["parser_used"] = "openpyxl"

        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            pg_num = sheet_idx + 1
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                rows.append(row_text)

            text = "\n".join(["\t".join(r) for r in rows])
            raw["_raw_text_by_page"][str(pg_num)] = text
            raw["_tables_by_page"][str(pg_num)] = [{"table_index": 0, "rows": rows}]
            raw["page_extractions"].append({
                "page_number": pg_num,
                "source_file": filepath.name,
                "extraction_method": "openpyxl",
                "best_text_chars": len(text)
            })
            print(f"  Sheet {pg_num} ({sheet_name}): {len(text)} chars, {len(rows)} rows")

        raw["_raw_evidence_status"] = "PRESENT"

    except Exception as e:
        raw["_raw_evidence_status"] = "UNREADABLE"
        raw["unreadable_log"] = {"reason": str(e)}
        print(f"  ERROR: {e}")

    return raw


def _extract_docx(filepath, raw):
    """Extract Word document."""
    try:
        from docx import Document
        doc = Document(str(filepath))
        raw["total_pages"] = 1
        raw["parser_used"] = "python-docx"

        text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        raw["_raw_text_by_page"]["1"] = text

        # Extract tables
        for t_idx, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            if "1" not in raw["_tables_by_page"]:
                raw["_tables_by_page"]["1"] = []
            raw["_tables_by_page"]["1"].append({"table_index": t_idx, "rows": rows})

        raw["page_extractions"].append({
            "page_number": 1,
            "source_file": filepath.name,
            "extraction_method": "python-docx",
            "best_text_chars": len(text)
        })
        raw["_raw_evidence_status"] = "PRESENT"
        print(f"  Document: {len(text)} chars")

    except Exception as e:
        raw["_raw_evidence_status"] = "UNREADABLE"
        raw["unreadable_log"] = {"reason": str(e)}
        print(f"  ERROR: {e}")

    return raw


def identify_ordering_pages(raw):
    """Phase 2: Scan raw evidence for ordering/code signals."""
    ordering_pages = []
    for pg, text in raw.get("_raw_text_by_page", {}).items():
        lower = text.lower()
        found_kw = [kw for kw in ORDERING_KEYWORDS if kw in lower]
        codes = re.findall(r'\b[A-Z]{2,5}[-]?[A-Z0-9]{2,}[-]?[0-9]{2,}[A-Z]?\b', text)
        real_codes = [c for c in codes if len(c) >= 5 and c not in ("INDIA", "MERIL", "ASTM", "CDSCO", "USFDA")]

        if found_kw or real_codes:
            ordering_pages.append({
                "page": pg,
                "keywords": found_kw,
                "codes": real_codes[:10],
                "has_table_signal": bool(found_kw)
            })

    return ordering_pages


def process_batch(start, end):
    """Process a batch of files with full evidence pipeline."""
    results = []

    for eid_num in range(start, end + 1):
        filename = get_file_for_eid(eid_num)
        if not filename:
            print(f"\n  SKIP {eid_num}: No file at this index")
            continue

        # Phase 1: Raw evidence extraction
        raw = extract_raw_evidence(eid_num)
        if not raw:
            continue

        # Phase 2: Identify ordering pages
        ordering = identify_ordering_pages(raw)

        eid = f"{eid_num:03d}"
        if ordering:
            print(f"\n  >> ORDERING SIGNALS FOUND on {len(ordering)} pages:")
            for op in ordering:
                print(f"     Page {op['page']}: keywords={op['keywords'][:3]}, codes={op['codes'][:5]}")
        else:
            print(f"\n  >> No ordering signals found. Marking for no-SKU verification.")

        results.append({
            "eid": eid,
            "filename": filename,
            "total_pages": raw.get("total_pages", 0),
            "evidence_status": raw.get("_raw_evidence_status", "?"),
            "ordering_pages": ordering,
            "has_ordering": len(ordering) > 0,
            "total_chars": sum(len(str(v)) for v in raw.get("_raw_text_by_page", {}).values())
        })

    # Summary
    print(f"\n{'='*100}")
    print(f"BATCH EXTRACTION SUMMARY (Files {start}-{end})")
    print(f"{'='*100}")

    with_ordering = [r for r in results if r["has_ordering"]]
    no_ordering = [r for r in results if not r["has_ordering"] and r["evidence_status"] == "PRESENT"]
    failed = [r for r in results if r["evidence_status"] != "PRESENT"]

    print(f"  Extracted: {len(results)}")
    print(f"  With ordering signals: {len(with_ordering)}")
    print(f"  No ordering signals: {len(no_ordering)}")
    print(f"  Failed/unreadable: {len(failed)}")

    if with_ordering:
        print(f"\n  FILES WITH ORDERING PAGES (need SKU extraction):")
        for r in with_ordering:
            pages = [op["page"] for op in r["ordering_pages"]]
            codes_sample = []
            for op in r["ordering_pages"]:
                codes_sample.extend(op["codes"][:3])
            print(f"    {r['eid']}: {r['filename'][:50]} | pages={pages} | sample_codes={codes_sample[:5]}")

    if no_ordering:
        print(f"\n  FILES WITH NO ORDERING SIGNALS (verified no-SKU):")
        for r in no_ordering:
            print(f"    {r['eid']}: {r['filename'][:50]} | {r['total_pages']} pages, {r['total_chars']} chars")

    if failed:
        print(f"\n  FAILED/UNREADABLE:")
        for r in failed:
            print(f"    {r['eid']}: {r['filename'][:50]} | status={r['evidence_status']}")

    # Save batch log
    LOG_DIR.mkdir(parents=True, exist_ok=True)
    log_path = LOG_DIR / f"batch3_extraction_log.json"
    with open(log_path, "w") as f:
        json.dump(results, f, indent=2)
    print(f"\n  Log saved: {log_path}")

    return results


if __name__ == "__main__":
    start = int(sys.argv[1]) if len(sys.argv) > 1 else 51
    end = int(sys.argv[2]) if len(sys.argv) > 2 else 75
    process_batch(start, end)
